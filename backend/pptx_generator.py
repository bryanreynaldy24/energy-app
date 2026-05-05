"""
PPTX Generator — Indonesia Energy Narrative Monitor
Produces the same 2-slide deck as the current manual process.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

def hex_color(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = {
    "navy":       "003878",
    "navyMid":    "005baa",
    "navyLight":  "e8f0f8",
    "red":        "c0392b",
    "redLight":   "fcecea",
    "amber":      "8b5e00",
    "amberLight": "fdf5e4",
    "green":      "1a6636",
    "greenLight": "e8f5ee",
    "ink":        "1a1a1a",
    "ink2":       "3c3c3c",
    "ink3":       "666666",
    "ink4":       "999999",
    "surface":    "f7f8fa",
    "surface2":   "f0f2f5",
    "rule":       "dde2e9",
    "white":      "FFFFFF",
    "purple":     "7c4dcc",
    "purpleLight":"f5f0fb",
}

def inches(v): return Inches(v)
def pt(v):     return Pt(v)

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width=None):
    shape = slide.shapes.add_shape(1, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(fill_hex)
    if line_hex:
        shape.line.color.rgb = hex_color(line_hex)
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_name="Calibri", font_size=6,
             bold=False, italic=False, color_hex="1a1a1a",
             align=PP_ALIGN.LEFT, valign="middle", wrap=True):
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_color(color_hex)
    return txBox

def gap_badge(gap):
    if gap > 5:
        return f"SM +{abs(gap):.1f}pt", C["red"], C["redLight"]
    elif gap < -5:
        return f"CM +{abs(gap):.1f}pt", C["amber"], C["amberLight"]
    else:
        lbl = f"SM +{abs(gap):.1f}pt" if gap > 0 else "≈ parity"
        return lbl, C["green"], C["greenLight"]

def generate_pptx(data: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = inches(13.3)
    prs.slide_height = inches(7.5)
    blank_layout = prs.slide_layouts[6]

    meta     = data["meta"]
    clusters = data["clusters"]
    regions  = data["regions"]
    daily_sm = data["daily_sm"]
    daily_cm = data["daily_cm"]
    satire   = data["satire_note"]

    # ── SLIDE 1: One-Pager ──────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)

    # Header
    add_rect(s1, 0, 0, 13.3, 0.65, C["navy"])
    add_text(s1, "INDONESIA ENERGY NARRATIVE MONITOR", 0.18, 0.05, 6, 0.3,
             "Georgia", 11, bold=True, color_hex=C["white"])
    add_text(s1, "One-Pager Analysis", 0.18, 0.32, 6, 0.22,
             "Georgia", 9, italic=True, color_hex="7db8e8")

    kpis = [
        (f"{meta['sm_total_posts']:,}", "SM POSTS", 9.0),
        (f"{meta['cm_total_articles']:,}", "CM ARTICLES", 10.5),
        (f"{meta['period_start']} – {meta['last_updated']}", "PERIOD", 11.9),
    ]
    for val, label, kx in kpis:
        add_text(s1, val,   kx, 0.04, 1.3, 0.28, "Georgia", 10, bold=True, color_hex=C["white"], align=PP_ALIGN.CENTER)
        add_text(s1, label, kx, 0.33, 1.3, 0.2,  "Courier New", 5.5, color_hex="7db8e8", align=PP_ALIGN.CENTER)

    # Column dividers
    add_rect(s1, 4.57, 0.75, 0.02, 6.4, C["rule"])
    add_rect(s1, 9.09, 0.75, 0.02, 6.4, C["rule"])

    # ── Col 1: Narrative Gap ──
    by = 0.75
    add_rect(s1, 0.13, by, 4.37, 0.2, C["navyLight"])
    add_text(s1, "NARRATIVE GAP ANALYSIS", 0.13, by, 4.37, 0.2,
             "Courier New", 6, bold=True, color_hex=C["navyMid"], align=PP_ALIGN.CENTER)
    by += 0.22
    add_text(s1, "Cluster Distribution & SM vs CM Divergence", 0.13, by, 4.37, 0.18,
             "Georgia", 7.5, bold=True, color_hex=C["navy"])
    by += 0.2

    # Table header
    add_rect(s1, 0.13, by, 4.37, 0.2, C["navy"])
    for txt, cx, cw, al in [
        ("CLUSTER", 0.13, 1.55, PP_ALIGN.LEFT),
        ("SM%",     1.68, 0.55, PP_ALIGN.CENTER),
        ("CM%",     2.23, 0.55, PP_ALIGN.CENTER),
        ("GAP",     2.78, 0.82, PP_ALIGN.CENTER),
        ("WHAT THIS CLUSTER IS ABOUT", 3.6,  0.9,  PP_ALIGN.LEFT),
    ]:
        add_text(s1, txt, cx+0.04, by, cw, 0.2, "Calibri", 5.5, bold=True,
                 color_hex=C["white"], align=al)
    by += 0.2

    row_bgs = [C["white"], C["surface"]]
    for i, cl in enumerate(clusters):
        rh = 0.72
        add_rect(s1, 0.13, by, 4.37, rh, row_bgs[i % 2], C["rule"], 0.3)
        accent = C["red"] if cl["gap"] > 3 else C["navyMid"] if cl["gap"] < -3 else C["green"]
        add_rect(s1, 0.13, by, 0.04, rh, accent)
        add_text(s1, cl["label"], 0.24, by+0.06, 1.4, 0.18, "Calibri", 6.2, bold=True, color_hex=C["ink"])
        add_text(s1, cl["sub"],   0.24, by+0.24, 1.4, 0.14, "Calibri", 5.0, color_hex=C["ink3"])
        sc = C["red"] if cl["gap"] > 0 else C["navyMid"]
        add_text(s1, f"{cl['sm_pct']}%", 1.68, by, 0.55, rh, "Georgia", 9, bold=True,
                 color_hex=sc, align=PP_ALIGN.CENTER)
        add_text(s1, f"{cl['cm_pct']}%", 2.23, by, 0.55, rh, "Georgia", 9,
                 color_hex=C["ink3"], align=PP_ALIGN.CENTER)
        glabel, gcol, gbg = gap_badge(cl["gap"])
        add_rect(s1, 2.82, by+0.22, 0.74, 0.22, gbg, gcol, 0.5)
        add_text(s1, glabel, 2.82, by+0.22, 0.74, 0.22, "Courier New", 5,
                 bold=True, color_hex=gcol, align=PP_ALIGN.CENTER)
        add_rect(s1, 3.54, by+0.08, 0.015, rh-0.16, C["rule"])
        add_text(s1, cl["desc"], 3.56, by+0.05, 0.9, rh-0.08, "Calibri", 4.9,
                 color_hex=C["ink2"])
        by += rh

    # Satire analyst note
    by += 0.08
    add_rect(s1, 0.13, by, 4.37, 0.72, C["purpleLight"], C["purple"], 0.6)
    add_rect(s1, 0.13, by, 0.04, 0.72, C["purple"])
    note = (f"ANALYST NOTE — G: Political Satire ({satire['sm_pct']}% SM, ~0% CM)  "
            "Not included as a main cluster due to small volume, but timing matters: satire spiked "
            "Apr 21–23 precisely between Broken Promise fading and Inflation Domino locking in. "
            "Content: memes, Bahlil mockery (\"King Bahlil\"), officials-not-solving-anything jokes. "
            "Acts as a transition narrative — anger that has cooled but keeps spreading. "
            "Near-zero CM coverage means conventional media is entirely missing this dimension.")
    add_text(s1, note, 0.23, by+0.05, 4.2, 0.63, "Calibri", 5.0, color_hex=C["ink2"])

    # ── Col 2: Regional Formation ──
    r2x, r2w = 4.65, 4.37
    ry = 0.75
    add_rect(s1, r2x, ry, r2w, 0.2, C["navyLight"])
    add_text(s1, "REGIONAL FORMATION", r2x, ry, r2w, 0.2,
             "Courier New", 6, bold=True, color_hex=C["navyMid"], align=PP_ALIGN.CENTER)
    ry += 0.22
    add_text(s1, "SM Engagement & Escalation Risk by Province", r2x, ry, r2w, 0.18,
             "Georgia", 7.5, bold=True, color_hex=C["navy"])
    ry += 0.2

    status_colors = {"hot": C["red"], "watch": C["amber"], "monitor": C["navyMid"]}
    status_bgs    = {"hot": "fef9f9", "watch": "fefcf5", "monitor": "f5f8fd"}
    status_labels = {"hot": "HOT 🔴", "watch": "WATCH 🟠", "monitor": "MONITOR 🔵"}
    card_h = (7.5 - ry - 0.35) / max(len(regions), 1) - 0.04

    for reg in regions:
        st = reg["status"]
        bc = status_colors.get(st, C["navyMid"])
        bg = status_bgs.get(st, "f5f8fd")
        add_rect(s1, r2x, ry, r2w, card_h, bg, bc, 0.8)
        add_rect(s1, r2x, ry, 0.06, card_h, bc)
        add_text(s1, reg["name"], r2x+0.1, ry+0.03, r2w-0.5, 0.18,
                 "Georgia", 6.5, bold=True, color_hex=C["ink"])
        pill_lbl = status_labels.get(st, st.upper())
        if reg["escalation"] > 0: pill_lbl += " ⚡"
        add_rect(s1, r2x+r2w-0.78, ry+0.04, 0.72, 0.17, status_bgs.get(st,"f5f8fd"), bc, 0.5)
        add_text(s1, pill_lbl, r2x+r2w-0.78, ry+0.04, 0.72, 0.17,
                 "Courier New", 4.8, bold=True, color_hex=bc, align=PP_ALIGN.CENTER)
        add_text(s1, f"Eng: {reg['engagement']:,.0f}", r2x+0.1, ry+0.21, 1.0, 0.14,
                 "Courier New", 5, color_hex=C["ink4"])
        dom_cl = next((c["name"] for c in clusters if c["id"] == reg["dominant_cluster"]), reg["dominant_cluster"])
        info = f"Dominant: {dom_cl} · {reg['sm_posts']} SM posts · {reg['cm_articles']} CM articles"
        if reg["escalation"] > 0:
            info += f" · {reg['escalation']} escalation signals"
        add_text(s1, info, r2x+0.1, ry+0.35, r2w-0.16, card_h-0.38, "Calibri", 5.2,
                 color_hex=C["ink2"])
        ry += card_h + 0.04

    # ── Col 3: Trend + Projection ──
    c3x, c3w = 9.17, 4.0
    t3y = 0.75
    add_rect(s1, c3x, t3y, c3w, 0.2, C["navyLight"])
    add_text(s1, "DAILY TREND + PROJECTION", c3x, t3y, c3w, 0.2,
             "Courier New", 6, bold=True, color_hex=C["navyMid"], align=PP_ALIGN.CENTER)
    t3y += 0.22
    add_text(s1, f"SM Engagement & CM Volume | {meta['period_start']} – {meta['last_updated']}",
             c3x, t3y, c3w, 0.18, "Georgia", 7.5, bold=True, color_hex=C["navy"])
    t3y += 0.22

    # Bar charts
    add_text(s1, "● SM ENGAGEMENT", c3x, t3y, c3w/2, 0.16, "Calibri", 5.5, color_hex=C["navyMid"], bold=True)
    add_text(s1, "● CM ARTICLES",   c3x+c3w/2+0.05, t3y, c3w/2, 0.16, "Calibri", 5.5, color_hex=C["navyMid"], bold=True)

    chart_h = 1.5
    chart_y_base = t3y + 0.18 + chart_h

    def draw_bars(slide, data_list, val_key, cx, cw, max_val, colors_fn):
        n = len(data_list)
        bw = (cw - 0.1) / n - 0.015
        for i, row in enumerate(data_list):
            v = row.get(val_key, 0) or 0
            bh = max(0.02, (v / max_val) * (chart_h - 0.15)) if max_val > 0 else 0.02
            bx = cx + 0.05 + i * (bw + 0.015)
            by2 = chart_y_base - bh
            color = colors_fn(i, v, data_list)
            add_rect(slide, bx, by2, bw, bh, color)
            lbl = row.get("date","")[-5:] if row.get("date") else ""
            add_text(slide, lbl, bx, chart_y_base+0.01, bw, 0.12,
                     "Calibri", 4, color_hex=C["ink4"], align=PP_ALIGN.CENTER)

    sm_vals = [r["engagement"] for r in daily_sm]
    cm_vals = [r["articles"]   for r in daily_cm]
    sm_max  = max(sm_vals) if sm_vals else 1
    cm_max  = max(cm_vals) if cm_vals else 1

    def sm_color(i, v, data):
        if v == max(r["engagement"] for r in data): return C["red"]
        if v > sm_max * 0.6: return C["amber"]
        return C["navyMid"]

    def cm_color(i, v, data):
        if v == max(r["articles"] for r in data): return C["red"]
        if v > cm_max * 0.6: return C["amber"]
        return C["navyMid"]

    draw_bars(s1, daily_sm, "engagement", c3x,          c3w/2-0.05, sm_max, sm_color)
    draw_bars(s1, daily_cm, "articles",   c3x+c3w/2+0.05, c3w/2-0.05, cm_max, cm_color)

    t3y = chart_y_base + 0.22

    # Observation box
    add_rect(s1, c3x, t3y, c3w, 0.55, C["surface"], C["rule"])
    add_rect(s1, c3x, t3y, 0.04, 0.55, C["navyMid"])
    add_text(s1,
             "TREND OBSERVATION  SM peaked driven by non-subsidized BBM price hike. "
             "Plateau persists signalling sustained public concern. CM peaked one day AHEAD of SM peak "
             "— media seeding preceded viral amplification. Secondary SM surge aligns with B50/cooking oil reports.",
             c3x+0.1, t3y+0.05, c3w-0.14, 0.47, "Calibri", 5.2, color_hex=C["ink2"])
    t3y += 0.63

    # Projections
    projections = [
        ("HIGH",   C["red"],   C["redLight"],
         "Subsidy Breach Escalation",
         "If Pertamax price adjustment confirmed, SM volume may spike 2–3× within 48h. B+C clusters co-activate.",
         "TRIGGER: Pertamax >IDR 14,000 OR official Pertalite revision announcement"),
        ("MEDIUM", C["amber"], C["amberLight"],
         "EV Narrative Breakout",
         "F-cluster SM surplus signals a fast-forming conversation alternatives haven't covered yet. Risk: misinformation.",
         "TRIGGER: EV incentive policy leaked OR Bahlil statement on kendaraan listrik"),
        ("LOWER",  C["green"], C["greenLight"],
         "Geopolitical De-escalation",
         "D-cluster (Iran/Hormuz) below 5% share. If oil prices stabilise, self-extinguishes without counter-messaging.",
         "TRIGGER: Brent crude < USD 75 sustained 3 days"),
    ]
    for level, lc, lb, name, desc, trigger in projections:
        ph = 0.6
        add_rect(s1, c3x, t3y, c3w, ph, C["surface"], C["rule"])
        add_rect(s1, c3x, t3y, 0.55, ph, lb, lc, 0.5)
        add_text(s1, level, c3x, t3y, 0.55, ph, "Courier New", 5.5, bold=True,
                 color_hex=lc, align=PP_ALIGN.CENTER)
        add_text(s1, name,    c3x+0.6, t3y+0.04, c3w-0.64, 0.16, "Calibri", 6.5, bold=True, color_hex=C["ink"])
        add_text(s1, desc,    c3x+0.6, t3y+0.2,  c3w-0.64, 0.24, "Calibri", 5.2, color_hex=C["ink2"])
        add_text(s1, trigger, c3x+0.6, t3y+0.45, c3w-0.64, 0.13, "Courier New", 4.8, color_hex=C["ink4"])
        t3y += ph + 0.05

    # Footer s1
    fy = 7.2
    add_rect(s1, 0, fy, 13.3, 0.3, C["surface2"], C["rule"])
    add_text(s1, "CONFIDENTIAL — Media Intelligence Internal Use",
             0.15, fy, 6, 0.3, "Courier New", 5.5, color_hex=C["ink4"])
    add_text(s1, f"Data: SM ({meta['sm_total_posts']:,} posts) + CM ({meta['cm_total_articles']:,} articles) | Keyword clustering | AVE-weighted | Last updated: {meta['last_updated']}",
             6.5, fy, 6.65, 0.3, "Courier New", 5.5, color_hex=C["ink4"], align=PP_ALIGN.RIGHT)

    # ── SLIDE 2: Source Verification ────────────────────────────────────────
    # (Placeholder — sources are dataset-specific and managed separately)
    s2 = prs.slides.add_slide(blank_layout)
    add_rect(s2, 0, 0, 13.3, 0.65, C["navy"])
    add_text(s2, "INDONESIA ENERGY NARRATIVE MONITOR", 0.18, 0.05, 6, 0.3,
             "Georgia", 11, bold=True, color_hex=C["white"])
    add_text(s2, "Source Verification — Every Claim, Linked", 0.18, 0.32, 6, 0.22,
             "Georgia", 9, italic=True, color_hex="7db8e8")
    add_text(s2, "Source verification links are maintained in the accompanying source log. "
             "All cluster assignments verified against post_translation (SM) and Title + Hit Sentence (CM).",
             0.5, 1.5, 12.3, 1.0, "Calibri", 9, color_hex=C["ink2"])

    prs.save(output_path)
